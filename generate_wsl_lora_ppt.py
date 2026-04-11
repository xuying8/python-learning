from pptx import Presentation
from pptx.util import Pt

OUTPUT_FILE = "advanced_wsl_lora_presentation.pptx"

prs = Presentation()


def add_bullet_slide(title: str, bullets: list[str], notes: str) -> None:
    slide_layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


# 1) Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Advanced WSL + LoRA Fine-Tuning Fundamentals"
slide.placeholders[1].text = "Practical engineering playbook"
slide.notes_slide.notes_text_frame.text = (
    "Set scope: production-oriented WSL workflows and efficient LLM adaptation with LoRA/QLoRA."
)

# Part 1: WSL
add_bullet_slide(
    "WSL Architecture (WSL1 vs WSL2)",
    [
        "WSL1: syscall translation; fast /mnt access",
        "WSL2: real Linux kernel in lightweight VM",
        "Use WSL2 for Docker, ML, kernel-level compatibility",
        "Choose per workload, not preference",
    ],
    "Explain why WSL2 is default for modern dev and ML. Mention compatibility vs IO tradeoffs.",
)

add_bullet_slide(
    "Performance Comparison (/mnt vs native Linux)",
    [
        "Native Linux FS (~/src) is faster for git/builds",
        "Mounted Windows paths (/mnt/c) are slower for metadata-heavy I/O",
        "Keep code + dependency caches inside ext4.vhdx",
        "Benchmark with real commands before migrating",
    ],
    "Show practical impact on pip install, npm, pytest, and large monorepos.",
)

add_bullet_slide(
    "Dev Workflow (VS Code + Docker)",
    [
        "Use VS Code Remote - WSL for Linux-native toolchains",
        "Enable Docker Desktop WSL integration per distro",
        "Run containers with source in Linux filesystem",
        "Use docker context and compose profiles per environment",
    ],
    "Walk through edit-build-run loop and common mount mistakes.",
)

add_bullet_slide(
    "GPU Usage in WSL",
    [
        "Requires WSL2 + compatible NVIDIA Windows driver",
        "Validate with nvidia-smi and framework checks",
        "Use CUDA/bf16 where supported for training speed",
        "Pin package versions to avoid CUDA mismatch",
    ],
    "Connect this to LoRA training: single-machine adapter fine-tuning workflows.",
)

add_bullet_slide(
    "WSL Best Practices",
    [
        "Enable systemd only for required services",
        "Tune .wslconfig memory/CPU/swap for workloads",
        "Automate port-forward refresh when WSL IP changes",
        "Export/import distros for reproducible environments",
    ],
    "Discuss operational hygiene, backups, and minimizing startup overhead.",
)

# Part 2: LoRA
add_bullet_slide(
    "What is LoRA?",
    [
        "Freeze base model; train small adapter matrices",
        "Weight update uses low-rank factorization ΔW = BA",
        "Far fewer trainable parameters than full fine-tuning",
        "Portable adapters per task/domain",
    ],
    "Clarify LoRA as parameter-efficient adaptation, not a new base model architecture.",
)

add_bullet_slide(
    "How LoRA Works (Diagram Explanation)",
    [
        "Original layer: y = Wx",
        "LoRA layer: y = Wx + (alpha/r) * BAx",
        "A and B are trainable; W stays frozen",
        "Adapter can be merged into base weights later",
    ],
    "Verbally describe flow: input x branches into base path and low-rank update path.",
)

add_bullet_slide(
    "LoRA Hyperparameters",
    [
        "Rank (r): adapter capacity vs memory",
        "Alpha: scaling strength of adapter update",
        "Learning rate: often higher than full fine-tuning",
        "Target modules: q/v (+ optional k/o/mlp)",
    ],
    "Provide starting defaults and explain how to tune based on loss/eval trends.",
)

add_bullet_slide(
    "Training Pipeline (HF PEFT)",
    [
        "Prepare clean instruction-response dataset",
        "Load base model/tokenizer and apply LoraConfig",
        "Train with Trainer/SFT loop + eval checkpoints",
        "Save adapter and inference config artifacts",
    ],
    "Mention deterministic preprocessing, validation split, and experiment tracking.",
)

add_bullet_slide(
    "QLoRA",
    [
        "Quantize base model to 4-bit (bitsandbytes)",
        "Train LoRA adapters in fp16/bf16",
        "Greatly reduces VRAM requirements",
        "Enables 7B/13B tuning on modest GPUs",
    ],
    "Explain tradeoffs: memory gains vs potential quality/throughput considerations.",
)

add_bullet_slide(
    "Real-World LoRA Use Cases",
    [
        "Domain QA assistants (legal, finance, support)",
        "Style/tone adaptation for enterprise copilots",
        "Task-specific classifiers via instruction tuning",
        "Multi-adapter serving per customer or product line",
    ],
    "Close with deployment choice: dynamic adapter loading vs merged models.",
)

prs.save(OUTPUT_FILE)
print(f"Created: {OUTPUT_FILE}")
